from pptx import Presentation

def process_pptx(file):
    presentation = Presentation(file)
    slides = []
    for slide in presentation.slides:
        slide_text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        slides.append(slide_text)
    return slides
