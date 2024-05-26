import PyPDF2
import docx
from pptx import Presentation
import nbformat
import nbconvert


def process_docx(file_path):
    doc = docx.Document(file_path)
    text = '\n'.join([para.text for para in doc.paragraphs])
    return text

def process_pptx(file_path):
    prs = Presentation(file_path)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

def process_pdf(file_path):
    text = ''
    with open(file_path, 'rb') as f:
        reader = PyPDF2.PdfFileReader(f)
        for page_num in range(reader.numPages):
            text += reader.getPage(page_num).extractText() + '\n'
    return text

def process_ipynb(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        notebook = nbformat.read(f, as_version=4)
        exporter = nbconvert.TextExporter()
        body, _ = exporter.from_notebook_node(notebook)
        return body
